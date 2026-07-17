"""Build the CableTwin Phase 2 final presentation (.pptx).

Implements the 9-slide specification from deck/outline.md with the CableTwin
palette, evidence screenshots and speaker notes. Output:
output/deck/CableTwin_SUPCOM_Final.pptx

Usage: python scripts/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "data-room" / "evidence" / "screenshots"
BRAND = ROOT / "brand"
OUT = ROOT / "output" / "deck" / "CableTwin_SUPCOM_Final.pptx"

DEEP = RGBColor(0x09, 0x29, 0x25)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
MINT = RGBColor(0xC9, 0xF3, 0x6B)
CREAM = RGBColor(0xF5, 0xF4, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE9, 0x76, 0x3B)
INK = RGBColor(0x10, 0x22, 0x1F)
SOFT = RGBColor(0x53, 0x64, 0x5F)
MIST = RGBColor(0x8F, 0xA8, 0xA2)

FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=None):
    sh = s.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    sh.fill.solid() if fill else sh.fill.background()
    if fill:
        sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = line_w or Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(4), wrap=True):
    """runs: list of paragraphs; each paragraph is a list of
    (txt, size, color, bold) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        for txt, size, color, bold in para:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def eyebrow(s, txt, color=TEAL, y=Inches(0.42)):
    text(s, Inches(0.7), y, Inches(12), Inches(0.35),
         [[(txt.upper(), 13, color, True)]])


def title(s, txt, color=INK, y=Inches(0.78), size=30):
    text(s, Inches(0.7), y, Inches(12), Inches(0.9), [[(txt, size, color, True)]])


def footer(s, dark=False):
    c = MIST if dark else SOFT
    text(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
         [[("CableTwin · SUPCOM · Oussama Akir · Fully synthetic data", 10, c, False)]])


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def pic(s, path, x, y, w=None, h=None):
    return s.shapes.add_picture(str(path), x, y, w, h)


def chip(s, x, y, w, txt, fill, fg, size=12, h=Inches(0.34)):
    c = box(s, x, y, w, h, fill=fill)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = fg
    r.font.bold = True
    return c


# ---------------------------------------------------------------- slide 1
s = slide(DEEP)
pic(s, BRAND / "logo" / "cabletwin-mark.png", Inches(0.7), Inches(0.55), h=Inches(1.05))
text(s, Inches(0.7), Inches(0.2), Inches(12), Inches(0.3),
     [[("AUTOMATE OR DIE 2026 · INDUSTRY · THEME 3", 12, MIST, True)]])
text(s, Inches(0.68), Inches(2.0), Inches(11.9), Inches(1.4),
     [[("Cable", 66, CREAM, True), ("Twin", 66, MINT, True)]])
text(s, Inches(0.7), Inches(3.35), Inches(11.2), Inches(1.5),
     [[("When a line stops, see the consequences", 32, CREAM, False)],
      [("before changing production.", 32, CREAM, False)]])
text(s, Inches(0.7), Inches(4.85), Inches(11), Inches(0.5),
     [[("The Waze of the cable factory", 20, MINT, False)]])
pic(s, SHOTS / "05-preview-service-gantt.png", Inches(8.6), Inches(5.35), w=Inches(4.2))
text(s, Inches(0.7), Inches(5.7), Inches(7.6), Inches(0.8),
     [[("Working offline prototype · Fictional Tunisian factory · Fully synthetic data",
        14, MIST, False)]])
footer(s, dark=True)
notes(s, "At 10:00, one production line stops. The schedule that looked correct one "
         "minute ago is now impossible. The production manager must decide which orders "
         "to move, delay or protect before every consequence is visible. CableTwin is the "
         "Waze of the cable factory: it shows what the disruption changes, compares "
         "feasible recovery routes and leaves the final decision to the manager. What you "
         "will see is a working offline prototype on a fully synthetic Tunisian cable factory.")

# ---------------------------------------------------------------- slide 2
s = slide()
eyebrow(s, "The problem")
title(s, "At 10:00, one stoppage exposes 3 orders and 620 minutes of delay", size=24)
text(s, Inches(0.7), Inches(1.7), Inches(6.6), Inches(1.6),
     [[("Which orders should be moved, delayed or protected —", 19, INK, False)],
      [("and what will each choice do to service, throughput", 19, INK, False)],
      [("and the workshop?", 19, INK, False)]])
labels = ["Eligible lines", "Durations", "Due times", "No overlap", "Planning horizon"]
x = Inches(0.7)
row_y = Inches(3.3)
for i, lab in enumerate(labels):
    if i == 3:  # second row
        x = Inches(0.7)
        row_y = Inches(3.78)
    w = Inches(0.28 + 0.115 * len(lab))
    chip(s, x, row_y, w, lab, DEEP, MINT, size=12)
    x = Emu(int(x) + int(w) + int(Inches(0.18)))
text(s, Inches(0.7), Inches(6.35), Inches(7.2), Inches(0.5),
     [[("The manager needs a feasible choice, not another dashboard.", 17, TEAL, True)]])
facts = [("Line 2 unavailable", "10:00 – 14:00"), ("Orders exposed", "3"),
         ("On time without adaptation", "7 / 10"), ("Projected cumulative delay", "620 min")]
fy = Inches(1.62)
for lab, val in facts:
    box(s, Inches(8.0), fy, Inches(4.6), Inches(1.06), fill=DEEP)
    text(s, Inches(8.25), Emu(int(fy) + int(Inches(0.13))), Inches(4.2), Inches(0.3),
         [[(lab.upper(), 11, MIST, True)]])
    text(s, Inches(8.25), Emu(int(fy) + int(Inches(0.4))), Inches(4.2), Inches(0.55),
         [[(val, 26, MINT if val != "10:00 – 14:00" else ORANGE, True)]])
    fy = Emu(int(fy) + int(Inches(1.22)))
text(s, Inches(0.7), Inches(4.5), Inches(6.9), Inches(1.7),
     [[("The active schedule becomes infeasible the moment the", 15, SOFT, False)],
      [("line stops: two orders overlap the stop, one downstream", 15, SOFT, False)],
      [("order is exposed, and every alternative must respect real", 15, SOFT, False)],
      [("constraints. The decision owner is the production manager.", 15, SOFT, False)]])
footer(s)
notes(s, "The process is incident-driven rescheduling. Once Line 2 is unavailable, two "
         "orders overlap the stop and one downstream order is exposed. Orders compete for "
         "the remaining eligible capacity. A valid answer must respect product-to-line "
         "compatibility, line-specific duration, downtime, non-overlap and the planning "
         "horizon. If every order stays on its original line, only seven of ten finish on "
         "time and cumulative delay reaches 620 minutes. The decision owner is the "
         "production manager; the process exists because the active schedule has become "
         "infeasible, while its local frequency and cost remain pilot questions.")

# ---------------------------------------------------------------- slide 3
s = slide()
eyebrow(s, "The product")
title(s, "CableTwin turns one disruption into three feasible routes")
steps = ["1  Detect the stop", "2  Reveal exposed orders and KPI",
         "3  Compare Service · Cost · Stability", "4  Preview, approve and record"]
sy = Inches(1.75)
for st in steps:
    box(s, Inches(0.7), sy, Inches(5.6), Inches(0.72), fill=DEEP)
    text(s, Inches(1.0), Emu(int(sy) + int(Inches(0.16))), Inches(5.2), Inches(0.4),
         [[(st, 16, CREAM, True)]])
    sy = Emu(int(sy) + int(Inches(0.92)))
box(s, Inches(6.9), Inches(1.75), Inches(5.7), Inches(2.4), fill=CREAM, line=TEAL, line_w=Pt(1.5))
text(s, Inches(7.2), Inches(2.0), Inches(5.1), Inches(1.9),
     [[("Like Waze after a road closes:", 17, INK, True)],
      [("same destination, different routes and visible", 15, SOFT, False)],
      [("trade-offs — the driver still decides.", 15, SOFT, False)]],
     space_after=Pt(6))
text(s, Inches(6.9), Inches(4.5), Inches(5.9), Inches(0.9),
     [[("One incident → three comparable recovery", 17, TEAL, True)],
      [("schedules → one traceable human decision", 17, TEAL, True)]])
box(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.66), fill=DEEP)
text(s, Inches(1.0), Inches(5.9), Inches(11.4), Inches(0.4),
     [[("Read-only decision support · No failure prediction · No machine control",
        15, MINT, True)]])
footer(s)
notes(s, "CableTwin is a lightweight incident-time digital twin. It keeps the minimum "
         "calculable state needed for this decision: lines, orders, calendars, "
         "compatibilities and the current schedule. When a route closes, it measures the "
         "impact and generates three alternatives. Service protects priority-weighted "
         "delivery, Cost minimizes the illustrative comparison cost, and Stability avoids "
         "line changes. The value is not a 3D factory. It is the ability to test a "
         "production decision safely before changing the real workshop.")

# ---------------------------------------------------------------- slide 4
s = slide()
eyebrow(s, "End-to-end workflow")
title(s, "One end-to-end loop connects the incident to an auditable decision", size=24)
cells = [("01 — Understand", "3 lines · 10/10 on time", "02-factory-nominal.png"),
         ("02 — Simulate", "Line 2 stops · 3 exposed", "03-incident-kpis.png"),
         ("03 — Compare", "3 feasible routes · same KPI", "04-strategies.png"),
         ("04 — Decide", "Manager approves · audit at 10:07", "07-approved-audit.png")]
cx = Inches(0.7)
cw = Inches(2.95)
for lab, cap, img in cells:
    text(s, cx, Inches(1.7), cw, Inches(0.3), [[(lab, 14, TEAL, True)]])
    pic(s, SHOTS / img, cx, Inches(2.1), w=cw)
    text(s, cx, Inches(3.95), cw, Inches(0.6), [[(cap, 12.5, INK, False)]])
    cx = Emu(int(cx) + int(cw) + int(Inches(0.12)))
box(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.66), fill=DEEP)
text(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.4),
     [[("The same executable path connects data, incident analysis, planning, "
        "Gantt preview, approval, audit and reset.", 13.5, CREAM, False)]])
footer(s)
notes(s, "This is the complete journey in the working build. We start from ten orders on "
         "time, trigger the four-hour stop, reveal the consequence, compare three plans, "
         "preview a revised Gantt and require the manager to approve. The final event "
         "records the choice at 10:07. Reset returns to the exact nominal state. These "
         "are post-freeze browser captures, not design mock-ups.")

# ---------------------------------------------------------------- slide 5
s = slide(DEEP)
eyebrow(s, "Working prototype · measured consequence", color=MINT)
text(s, Inches(0.7), Inches(0.78), Inches(12), Inches(0.9),
     [[("Service cuts projected delay from ", 30, CREAM, True),
       ("620 to 140 minutes", 30, MINT, True)]])
mets = [("620 → 140 min", "projected cumulative delay"),
        ("188 → 216 km", "completed by 18:00 · +14.9%"),
        ("180 → 30 min", "overtime · −83.3%"),
        ("7/10 → 8/10", "orders on time · 3 line moves")]
mx = Inches(0.7)
for val, lab in mets:
    text(s, mx, Inches(1.75), Inches(3.0), Inches(0.5), [[(val, 21, MINT, True)]])
    text(s, mx, Inches(2.2), Inches(3.0), Inches(0.5), [[(lab, 12, MIST, False)]])
    mx = Emu(int(mx) + int(Inches(3.05)))
p = pic(s, SHOTS / "05-preview-service-gantt.png", Inches(0.7), Inches(2.95),
        w=Inches(8.3), h=Inches(2.96))
p.crop_top = 0.305
p.crop_bottom = 0.06
rows = [("Stability", "620 min · 188 km · 0 moves"),
        ("Service", "140 min · 216 km · 3 moves"),
        ("Cost", "170 min · 216 km · 2 moves")]
ry = Inches(2.9)
for name, vals in rows:
    fill = TEAL if name == "Service" else DEEP
    b = box(s, Inches(9.25), ry, Inches(3.35), Inches(0.92), fill=fill,
            line=TEAL if name != "Service" else None, line_w=Pt(1))
    text(s, Inches(9.5), Emu(int(ry) + int(Inches(0.1))), Inches(3.0), Inches(0.3),
         [[(name, 14, MINT if name == "Service" else CREAM, True)]])
    text(s, Inches(9.5), Emu(int(ry) + int(Inches(0.42))), Inches(3.0), Inches(0.4),
         [[(vals, 13, CREAM, False)]])
    ry = Emu(int(ry) + int(Inches(1.08)))
text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
     [[("Fixed synthetic scenario · calculated by the working prototype · not factory ROI",
        12.5, MIST, False)]])
footer(s, dark=True)
notes(s, "Here is the live journey. The nominal schedule starts at ten of ten. I trigger "
         "the stop on Line 2. Without adaptation, three orders are exposed, delay reaches "
         "620 minutes, overtime reaches 180 minutes and only 188 kilometres finish by "
         "18:00. CableTwin calculates three routes from the same incident. Service reduces "
         "delay to 140 minutes and overtime to 30, while output reaches 216 kilometres. It "
         "accepts three line moves to protect priority-weighted delivery. Cost uses two "
         "moves and reaches 170 minutes of delay. Stability makes no move and keeps the "
         "full 620 minutes. I preview Service, inspect the revised Gantt and approve. The "
         "audit event appears; no order is sent to a machine. "
         "[Phase 3: run the live 90-second journey here; if the environment fails, "
         "narrate these frozen captures.]")

# ---------------------------------------------------------------- slide 6
s = slide()
eyebrow(s, "Core AI")
title(s, "The AI is the planning decision, not the dashboard")
box(s, Inches(0.7), Inches(1.7), Inches(5.75), Inches(3.5), fill=DEEP)
text(s, Inches(1.0), Inches(1.95), Inches(5.2), Inches(3.1),
     [[("LIVE SYMBOLIC PLANNER", 13, MINT, True)],
      [("Decision variables", 13, MIST, True)],
      [("line assignment · sequence · start / end", 13.5, CREAM, False)],
      [("Hard constraints", 13, MIST, True)],
      [("eligibility · duration · downtime · no overlap · horizon", 13.5, CREAM, False)],
      [("Business policies", 13, MIST, True)],
      [("Service · Cost · Stability (lexicographic objectives)", 13.5, CREAM, False)]],
     space_after=Pt(7))
box(s, Inches(6.85), Inches(1.7), Inches(5.75), Inches(3.5), fill=CREAM, line=TEAL, line_w=Pt(1.5))
text(s, Inches(7.15), Inches(1.95), Inches(5.2), Inches(3.1),
     [[("INDEPENDENT BOUNDED ORACLE", 13, TEAL, True)],
      [("17,856 candidate schedules evaluated", 15, INK, True)],
      [("10,440 feasible complete schedules", 15, INK, True)],
      [("3/3 displayed routes = unique policy optimum", 15, INK, True)],
      [("9/9 automated checks pass", 15, INK, True)]],
     space_after=Pt(9))
box(s, Inches(0.7), Inches(5.45), Inches(11.9), Inches(1.0), fill=DEEP)
text(s, Inches(1.0), Inches(5.6), Inches(11.4), Inches(0.75),
     [[("A 9/10 route exists, but creates 230 min of total delay; Service chooses 8/10 "
        "with only 140 min", 14, CREAM, False)],
      [("because priority-weighted lateness matters. Exact for the encoded bounded demo — "
        "not factory-scale optimization.", 14, MIST, False)]])
footer(s)
notes(s, "The intelligence is symbolic planning under constraints. For seven orders not "
         "complete at incident time, the engine chooses eligible lines, sequences and "
         "times. A route is rejected if it breaks compatibility, duration, downtime, "
         "non-overlap or the horizon. Service, Cost and Stability then express explicit "
         "lexicographic objectives. The browser uses a fast deterministic constructive "
         "planner. Separately, an exhaustive oracle enumerates 17,856 assignment and "
         "sequence candidates, finds 10,440 feasible schedules and confirms that all "
         "three displayed plans are the unique optima of their declared policies in this "
         "bounded scenario. That separation gives us both a usable product and an honest "
         "certificate.")

# ---------------------------------------------------------------- slide 7
s = slide()
eyebrow(s, "From simulation to industrial proof")
title(s, "Simulation proves the mechanism; a shadow pilot will prove value", size=24)
box(s, Inches(0.7), Inches(1.7), Inches(3.9), Inches(4.3), fill=DEEP)
text(s, Inches(1.0), Inches(1.95), Inches(3.4), Inches(3.8),
     [[("PROVED IN DEMO", 14, MINT, True)],
      [("Same incident", 14, CREAM, False)],
      [("Same KPI definitions", 14, CREAM, False)],
      [("Reproducible before/after", 14, CREAM, False)],
      [("calculations", 14, CREAM, False)]], space_after=Pt(8))
steps7 = ["5–10 historical stoppages", "Same information cutoff as the real decision",
          "Planner-validated constraints", "CableTwin vs actual historical baseline",
          "Time, delay, output, overtime, moves, usability"]
sy = Inches(1.7)
for st in steps7:
    box(s, Inches(4.95), sy, Inches(7.65), Inches(0.62), fill=CREAM, line=TEAL, line_w=Pt(1))
    text(s, Inches(5.2), Emu(int(sy) + int(Inches(0.13))), Inches(7.2), Inches(0.35),
         [[(st, 14, INK, False)]])
    sy = Emu(int(sy) + int(Inches(0.78)))
gates = "30% faster assessment · 20% lower delay · 0 hard-constraint violations · ≥70% planner usability · 100% traceability"
box(s, Inches(4.95), Inches(5.75), Inches(7.65), Inches(0.85), fill=DEEP)
text(s, Inches(5.2), Inches(5.88), Inches(7.2), Inches(0.6),
     [[("PROPOSED SUCCESS GATES — targets to agree, not achieved claims", 11.5, MINT, True)],
      [(gates, 12, CREAM, False)]])
text(s, Inches(0.7), Inches(6.3), Inches(3.9), Inches(0.5),
     [[("Ends in Go · Revise · Stop", 14, TEAL, True)]])
footer(s)
notes(s, "I do not convert synthetic results into a factory ROI. Phase 2 proves that the "
         "workflow calculates meaningful, reproducible trade-offs. A real proof starts "
         "with five to ten past stoppages. We freeze the information available at the "
         "decision time, keep the actual decision as baseline, validate the rules with "
         "the planner and compare identical KPI definitions. Proposed gates are 30 "
         "percent faster assessment, 20 percent less delay, zero hard-constraint "
         "violations, at least 70 percent planner usability and complete traceability. "
         "The sponsor decides go, revise or stop.")

# ---------------------------------------------------------------- slide 8
s = slide()
eyebrow(s, "Business viability")
title(s, "Start read-only in one workshop; scale only after measured proof", size=24)
text(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.4),
     [[("User: ", 16, SOFT, False), ("Production manager / planner", 16, INK, True),
       ("      Buyer: ", 16, SOFT, False), ("Plant or operations director", 16, INK, True),
       ("      Pilot: ", 16, SOFT, False), ("4–6 weeks · read-only", 16, INK, True)]])
path8 = ["Validate process", "Map read-only data", "Encode workshop constraints",
         "Replay and shadow", "Measured go / revise / stop"]
px = Inches(0.7)
pw = Inches(2.32)
for i, st in enumerate(path8):
    fill = DEEP if i < 4 else TEAL
    box(s, px, Inches(2.4), pw, Inches(0.95), fill=fill)
    text(s, Emu(int(px) + int(Inches(0.15))), Inches(2.55), Emu(int(pw) - int(Inches(0.3))),
         Inches(0.7), [[(st, 13, CREAM if i < 4 else MINT, True)]])
    px = Emu(int(px) + int(pw) + int(Inches(0.1)))
comp = [("Manual / spreadsheet", "familiar but hard to compare and trace"),
        ("ERP · MES · APS", "system of record / broad planning"),
        ("High-fidelity twin", "rich but integration-heavy"),
        ("CableTwin", "focused incident layer · 3 explainable routes · human approval")]
cy = Inches(3.85)
for name, desc in comp:
    is_ct = name == "CableTwin"
    box(s, Inches(0.7), cy, Inches(11.9), Inches(0.52),
        fill=TEAL if is_ct else CREAM, line=None if is_ct else TEAL, line_w=Pt(0.75))
    text(s, Inches(0.95), Emu(int(cy) + int(Inches(0.1))), Inches(3.4), Inches(0.3),
         [[(name, 13, MINT if is_ct else TEAL, True)]])
    text(s, Inches(4.4), Emu(int(cy) + int(Inches(0.1))), Inches(8.0), Inches(0.3),
         [[(desc, 13, CREAM if is_ct else SOFT, False)]])
    cy = Emu(int(cy) + int(Inches(0.62)))
box(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5), fill=DEEP)
text(s, Inches(1.0), Inches(6.55), Inches(11.4), Inches(0.3),
     [[("No prediction · no autonomous control · no write-back · no price or ROI claim",
        13, MINT, True)]])
footer(s)
notes(s, "The primary user is the production manager; the likely sponsor is the plant or "
         "operations director. The first engagement is deliberately small: one workshop, "
         "four to six weeks, read-only exports and five to ten historical incidents. It "
         "requires no sensors, GPU, cloud migration or machine access. CableTwin "
         "complements the systems already present. A spreadsheet remains flexible, an ERP "
         "or MES remains the system of record, an APS may already solve broad planning, "
         "and a high-fidelity twin models much more. CableTwin earns its place only where "
         "incident response is still manual, difficult to compare or weakly traced. If "
         "constraints, data or KPI value fail the gate, the pilot stops.")

# ---------------------------------------------------------------- slide 9
s = slide(DEEP)
eyebrow(s, "Ready for phase 3", color=MINT)
text(s, Inches(0.7), Inches(0.78), Inches(12), Inches(0.8),
     [[("CableTwin is ready for Phase 3", 32, CREAM, True)]])
grid = [("25% · Business process & impact",
         "One incident-rescheduling loop · delay, overtime and output measured"),
        ("20% · Core AI",
         "Constraint planning · 17,856-candidate independent bounded certificate"),
        ("20% · Working prototype",
         "Offline incident → routes → Gantt → approval → audit → reset · 9/9 checks"),
        ("20% · Business viability",
         "Named user and buyer · 4–6-week read-only pilot · go/no-go gates"),
        ("15% · Pitch & clarity",
         "One Waze metaphor · three routes · one human decision")]
gy = Inches(1.8)
for crit, ev in grid:
    box(s, Inches(0.7), gy, Inches(3.85), Inches(0.68), fill=TEAL)
    text(s, Inches(0.95), Emu(int(gy) + int(Inches(0.16))), Inches(3.5), Inches(0.4),
         [[(crit, 13.5, CREAM, True)]])
    text(s, Inches(4.85), Emu(int(gy) + int(Inches(0.16))), Inches(7.8), Inches(0.4),
         [[(ev, 13.5, CREAM, False)]])
    gy = Emu(int(gy) + int(Inches(0.84)))
box(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.72), fill=MINT)
text(s, Inches(1.0), Inches(6.36), Inches(11.4), Inches(0.4),
     [[("Select CableTwin for the final pitch and see the complete live decision "
        "in 90 seconds.", 16, DEEP, True)]])
footer(s, dark=True)
notes(s, "CableTwin is one coherent proof across the official grid. The process is a "
         "specific line-stoppage rescheduling decision. Impact is calculated with "
         "explicit definitions. The core AI builds constrained alternatives and is "
         "independently verified. The prototype runs end to end, offline, with nine "
         "automated checks. The business path begins with a read-only pilot and stops if "
         "evidence is not there. And the story remains simple: when a production route "
         "closes, show three feasible routes before the manager decides. Select CableTwin "
         "for Phase 3, and I will demonstrate that complete decision live in 90 seconds.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"deck written: {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
